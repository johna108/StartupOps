from fastapi import FastAPI, APIRouter, Depends, HTTPException, Request
from fastapi.responses import FileResponse, StreamingResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
import os
import logging
import uuid
import io
import json
import re
from pathlib import Path
from pydantic import BaseModel
from typing import List, Optional
from datetime import datetime, timezone
from pptx import Presentation
from pptx.util import Inches, Pt
from supabase import create_client, Client
from functools import lru_cache
import time

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# Supabase client for auth verification AND database
supabase_url = os.environ.get('SUPABASE_URL')
supabase_service_key = os.environ.get('SUPABASE_SERVICE_ROLE_KEY')
if not supabase_url or not supabase_service_key:
    raise ValueError("SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY must be set in .env")

supabase_client: Client = create_client(supabase_url, supabase_service_key)

# Gemini API key
GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY')

app = FastAPI()
api_router = APIRouter(prefix="/api")

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Simple in-memory cache for auth tokens (limits external calls to Supabase users endpoint)
# Format: {token: (user_object, timestamp)}
_token_cache = {}

def get_cached_user(token: str):
    now = time.time()
    # expired after 60 seconds
    if token in _token_cache:
        user, timestamp = _token_cache[token]
        if now - timestamp < 60:
            return user
    return None

def set_cached_user(token: str, user):
    _token_cache[token] = (user, time.time())
    # clear old cache if too big
    if len(_token_cache) > 1000:
        _token_cache.clear()

# ==================== CORS CONFIGURATION ====================

cors_origins_env = os.environ.get('CORS_ORIGINS', '*')
default_origins = [
    "https://startup-gamma-seven.vercel.app",
    "http://localhost:3000",
    "http://127.0.0.1:3000",
    "http://localhost:8001",
    "http://localhost:8002",
]

if cors_origins_env == '*':
    cors_origins = ['*']
else:
    cors_origins = [o.strip() for o in cors_origins_env.split(',') if o.strip()]
    for default in default_origins:
        if default not in cors_origins:
            cors_origins.append(default)

logger.info(f"CORS configured: {cors_origins}")

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=cors_origins if cors_origins else ['*'],
    allow_methods=["*"],
    allow_headers=["*"],
    allow_origin_regex=".*",
)

# Handle OPTIONS requests (preflight)
@app.options("/{full_path:path}")
async def preflight_handler(full_path: str):
    return {"status": "ok"}

# ==================== PYDANTIC MODELS ====================

class ProfileCreate(BaseModel):
    full_name: Optional[str] = None

class StartupCreate(BaseModel):
    name: str
    description: Optional[str] = None
    industry: Optional[str] = None
    stage: Optional[str] = "idea"
    website: Optional[str] = None
    initial_role: Optional[str] = "founder"

class StartupUpdate(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    industry: Optional[str] = None
    stage: Optional[str] = None
    website: Optional[str] = None

class TaskCreate(BaseModel):
    title: str
    description: Optional[str] = None
    status: Optional[str] = "todo"
    priority: Optional[str] = "medium"
    assigned_to: Optional[str] = None
    milestone_id: Optional[str] = None
    due_date: Optional[str] = None

class TaskUpdate(BaseModel):
    title: Optional[str] = None
    description: Optional[str] = None
    status: Optional[str] = None
    priority: Optional[str] = None
    assigned_to: Optional[str] = None
    milestone_id: Optional[str] = None
    due_date: Optional[str] = None

class MilestoneCreate(BaseModel):
    title: str
    description: Optional[str] = None
    target_date: Optional[str] = None

class MilestoneUpdate(BaseModel):
    title: Optional[str] = None
    description: Optional[str] = None
    target_date: Optional[str] = None
    status: Optional[str] = None

class FeedbackCreate(BaseModel):
    title: str
    content: Optional[str] = None
    category: Optional[str] = "product"
    rating: Optional[int] = 3
    source: Optional[str] = "internal"

class AIInsightRequest(BaseModel):
    startup_id: str
    prompt_type: Optional[str] = "general"
    custom_prompt: Optional[str] = None

class PitchRequest(BaseModel):
    startup_id: str
    custom_prompt: Optional[str] = None

class SubscriptionUpdate(BaseModel):
    plan: str

class MemberRoleUpdate(BaseModel):
    role: str

class TaskStatusUpdate(BaseModel):
    status: str

class IncomeCreate(BaseModel):
    title: str
    amount: float
    category: Optional[str] = "revenue"
    date: Optional[str] = None
    notes: Optional[str] = None

class ExpenseCreate(BaseModel):
    title: str
    amount: float
    category: Optional[str] = "operations"
    date: Optional[str] = None
    notes: Optional[str] = None

class InvestmentCreate(BaseModel):
    investor_name: str
    amount: float
    equity_percentage: Optional[float] = 0
    investment_type: Optional[str] = "seed"
    date: Optional[str] = None
    notes: Optional[str] = None

class JoinStartupRequest(BaseModel):
    invite_code: str

class SignupRequest(BaseModel):
    email: str
    password: str
    full_name: Optional[str] = None

# ==================== AUTH DEPENDENCY ====================

async def get_current_user(request: Request):
    auth_header = request.headers.get('authorization', '')
    if not auth_header.startswith('Bearer '):
        logger.warning("Missing or invalid Authorization header format")
        raise HTTPException(status_code=401, detail="Not authenticated")
    token = auth_header.split(' ')[1]
    
    # Try cache first
    cached_user = get_cached_user(token)
    if cached_user:
        return cached_user

    if not token or token == 'undefined' or token == 'null':
        logger.warning(f"Invalid token value: {token[:20] if token else 'empty'}...")
        raise HTTPException(status_code=401, detail="No valid token provided")
    try:
        user_response = supabase_client.auth.get_user(token)
        if not user_response or not user_response.user:
            logger.error("Supabase returned no user for token")
            raise HTTPException(status_code=401, detail="Invalid or expired token")
        
        # Cache the result
        set_cached_user(token, user_response.user)
        
        logger.info(f"User authenticated: {user_response.user.email}")
        return user_response.user
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Auth error: {type(e).__name__}: {e}")
        raise HTTPException(status_code=401, detail="Invalid or expired token")

# ==================== HEALTH CHECK ====================

@api_router.get("/")
async def root():
    return {"message": "StartupOps API is running (Supabase PostgreSQL)"}

# ==================== AUTH ROUTES ====================

@api_router.post("/auth/signup")
async def signup_user(body: SignupRequest):
    """Create user with auto-confirm."""
    try:
        user_response = supabase_client.auth.admin.create_user({
            "email": body.email,
            "password": body.password,
            "email_confirm": True,
            "user_metadata": {"full_name": body.full_name or body.email.split("@")[0]},
        })
        return {"message": "Account created", "user_id": user_response.user.id}
    except Exception as e:
        error_msg = str(e)
        if "already been registered" in error_msg or "already exists" in error_msg:
            raise HTTPException(status_code=409, detail="An account with this email already exists.")
        logger.error(f"Signup error: {e}")
        raise HTTPException(status_code=400, detail="Failed to create account")

@api_router.post("/auth/verify")
async def verify_and_sync_profile(body: ProfileCreate, user=Depends(get_current_user)):
    response = supabase_client.table("profiles").select("*").eq("id", str(user.id)).execute()
    existing = response.data[0] if response.data else None
    
    if existing:
        return existing
        
    profile = {
        "id": str(user.id),
        "email": user.email,
        "full_name": body.full_name or user.user_metadata.get("full_name", "") or user.email.split("@")[0],
        "avatar_url": user.user_metadata.get("avatar_url", ""),
        "created_at": datetime.now(timezone.utc).isoformat(),
        "updated_at": datetime.now(timezone.utc).isoformat(),
    }
    supabase_client.table("profiles").insert(profile).execute()
    return profile

@api_router.get("/auth/me")
async def get_me(user=Depends(get_current_user)):
    response = supabase_client.table("profiles").select("*").eq("id", str(user.id)).execute()
    profile = response.data[0] if response.data else None
    
    if not profile:
        profile = {
            "id": str(user.id),
            "email": user.email,
            "full_name": user.user_metadata.get("full_name", "") or user.email.split("@")[0],
            "avatar_url": user.user_metadata.get("avatar_url", ""),
            "created_at": datetime.now(timezone.utc).isoformat(),
            "updated_at": datetime.now(timezone.utc).isoformat(),
        }
        supabase_client.table("profiles").insert(profile).execute()
    return profile

@api_router.put("/auth/profile")
async def update_profile(body: ProfileCreate, user=Depends(get_current_user)):
    updates = {"updated_at": datetime.now(timezone.utc).isoformat()}
    if body.full_name:
        updates["full_name"] = body.full_name
    
    response = supabase_client.table("profiles").update(updates).eq("id", str(user.id)).execute()
    return response.data[0] if response.data else None

# ==================== STARTUP ROUTES ====================

@api_router.post("/startups")
async def create_startup(body: StartupCreate, user=Depends(get_current_user)):
    try:
        # Check if creating investor profile
        if body.initial_role == "investor":
            # For investors, we create a special "My Portfolio" startup
            # This is a bit of a hack to fit the existing "currentStartup" context model
            startup_name = body.name or f"{user.user_metadata.get('full_name', 'My')} Portfolio"
            description = "Investor Portfolio"
            industry = "Investment"
            stage = "active"
        else:
            startup_name = body.name
            description = body.description
            industry = body.industry
            stage = body.stage

        invite_code = str(uuid.uuid4()).replace("-", "")[:8].upper()
        
        # 1. Create Startup Entry
        startup_data = {
            "name": startup_name,
            "description": description,
            "industry": industry,
            "stage": stage,
            "website": body.website,
            "founder_id": str(user.id),
            "invite_code": invite_code,
            "subscription_plan": "free" if body.initial_role != "investor" else "investor_pro",
            "created_at": datetime.now(timezone.utc).isoformat(),
            "updated_at": datetime.now(timezone.utc).isoformat(),
        }
        
        s_res = supabase_client.table("startups").insert(startup_data).execute()
        if not s_res.data:
             raise HTTPException(status_code=500, detail="Failed to create startup record")
             
        new_startup = s_res.data[0]
        
        # 2. Add User as Member with correct role
        member_data = {
            "startup_id": new_startup['id'],
            "user_id": str(user.id),
            "role": body.initial_role or "founder",  # Use the requested role
            "joined_at": datetime.now(timezone.utc).isoformat(),
        }
        
        m_res = supabase_client.table("startup_members").insert(member_data).execute()
        
        # Return with role attached for frontend convenience
        new_startup['user_role'] = member_data['role']
        return new_startup

    except Exception as e:
        logger.error(f"Create startup error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@api_router.get("/startups")
async def get_user_startups(user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("user_id", str(user.id)).execute()
    memberships = mem_response.data
    
    if not memberships:
        return []
        
    startup_ids = [m["startup_id"] for m in memberships]
    startups_response = supabase_client.table("startups").select("*").in_("id", startup_ids).execute()
    startups = startups_response.data
    
    for s in startups:
        membership = next((m for m in memberships if m["startup_id"] == s["id"]), None)
        s["user_role"] = membership["role"] if membership else "member"
    
    return startups

@api_router.get("/startups/{startup_id}")
async def get_startup(startup_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    member = mem_response.data[0] if mem_response.data else None
    
    if not member:
        raise HTTPException(status_code=403, detail="Not a member of this startup")
        
    startup_response = supabase_client.table("startups").select("*").eq("id", startup_id).execute()
    startup = startup_response.data[0] if startup_response.data else None
    
    if not startup:
        raise HTTPException(status_code=404, detail="Startup not found")
        
    startup["user_role"] = member["role"]
    return startup

@api_router.put("/startups/{startup_id}")
async def update_startup(startup_id: str, body: StartupUpdate, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    member = mem_response.data[0] if mem_response.data else None
    
    if not member or member["role"] != "founder":
        raise HTTPException(status_code=403, detail="Only founders can update startup")
        
    updates = {"updated_at": datetime.now(timezone.utc).isoformat()}
    for field in ["name", "description", "industry", "stage", "website"]:
        val = getattr(body, field, None)
        if val is not None:
            updates[field] = val
            
    response = supabase_client.table("startups").update(updates).eq("id", startup_id).execute()
    return response.data[0] if response.data else None

@api_router.post("/startups/join")
async def join_startup(body: JoinStartupRequest, user=Depends(get_current_user)):
    startup_response = supabase_client.table("startups").select("*").eq("invite_code", body.invite_code).execute()
    startup = startup_response.data[0] if startup_response.data else None
    
    if not startup:
        raise HTTPException(status_code=404, detail="Invalid invite code")
        
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup["id"]).eq("user_id", str(user.id)).execute()
    existing = mem_response.data[0] if mem_response.data else None
    
    if existing:
        raise HTTPException(status_code=400, detail="Already a member")
        
    count_response = supabase_client.table("startup_members").select("id", count="exact").eq("startup_id", startup["id"]).execute()
    member_count = count_response.count
    
    plan = startup.get("subscription_plan", "free")
    max_members = 5 if plan == "free" else 999
    
    if member_count >= max_members:
        raise HTTPException(status_code=400, detail=f"Team limit reached for {plan} plan")
        
    member = {
        "id": str(uuid.uuid4()),
        "startup_id": startup["id"],
        "user_id": str(user.id),
        "role": "member",
        "joined_at": datetime.now(timezone.utc).isoformat(),
    }
    supabase_client.table("startup_members").insert(member).execute()
    return startup

# ==================== TASK ROUTES ====================

@api_router.post("/startups/{startup_id}/tasks")
async def create_task(startup_id: str, body: TaskCreate, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    task = {
        "id": str(uuid.uuid4()),
        "startup_id": startup_id,
        "title": body.title,
        "description": body.description or "",
        "status": body.status or "todo",
        "priority": body.priority or "medium",
        "assigned_to": body.assigned_to if body.assigned_to else None,
        "created_by": str(user.id),
        "milestone_id": body.milestone_id if body.milestone_id else None,
        "due_date": body.due_date if body.due_date else None,
        "created_at": datetime.now(timezone.utc).isoformat(),
        "updated_at": datetime.now(timezone.utc).isoformat(),
    }
    
    supabase_client.table("tasks").insert(task).execute()
    return task

@api_router.get("/startups/{startup_id}/tasks")
async def get_tasks(startup_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    tasks_response = supabase_client.table("tasks").select("*").eq("startup_id", startup_id).execute()
    return tasks_response.data

@api_router.put("/tasks/{task_id}")
async def update_task(task_id: str, body: TaskUpdate, user=Depends(get_current_user)):
    task_response = supabase_client.table("tasks").select("*").eq("id", task_id).execute()
    task = task_response.data[0] if task_response.data else None
    
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
        
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", task["startup_id"]).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    updates = {"updated_at": datetime.now(timezone.utc).isoformat()}
    for field in ["title", "description", "status", "priority", "assigned_to", "milestone_id", "due_date"]:
        val = getattr(body, field, None)
        # Handle empty strings for UUID/Date fields which cause Postgres errors
        if field in ["assigned_to", "milestone_id", "due_date"] and val == "":
            updates[field] = None
        elif val is not None:
            updates[field] = val
            
    response = supabase_client.table("tasks").update(updates).eq("id", task_id).execute()
    return response.data[0] if response.data else None

@api_router.delete("/tasks/{task_id}")
async def delete_task(task_id: str, user=Depends(get_current_user)):
    task_response = supabase_client.table("tasks").select("*").eq("id", task_id).execute()
    task = task_response.data[0] if task_response.data else None
    
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
        
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", task["startup_id"]).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    supabase_client.table("tasks").delete().eq("id", task_id).execute()
    return {"success": True}

@api_router.patch("/tasks/{task_id}/status")
async def update_task_status(task_id: str, body: TaskStatusUpdate, user=Depends(get_current_user)):
    task_response = supabase_client.table("tasks").select("*").eq("id", task_id).execute()
    task = task_response.data[0] if task_response.data else None
    
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", task["startup_id"]).eq("user_id", str(user.id)).execute()
    member = mem_response.data[0] if mem_response.data else None
    
    if not member:
        raise HTTPException(status_code=403, detail="Not a member")
    
    is_assigned = task.get("assigned_to") == str(user.id)
    is_manager_or_founder = member["role"] in ["founder", "manager"]
    
    if not is_assigned and not is_manager_or_founder:
        raise HTTPException(status_code=403, detail="You can only update status of tasks assigned to you")
    
    valid_statuses = ["todo", "in_progress", "review", "done"]
    if body.status not in valid_statuses:
        raise HTTPException(status_code=400, detail=f"Invalid status. Must be one of: {valid_statuses}")
    
    response = supabase_client.table("tasks").update(
        {"status": body.status, "updated_at": datetime.now(timezone.utc).isoformat()}
    ).eq("id", task_id).execute()
    
    return response.data[0] if response.data else None

# ==================== MILESTONE ROUTES ====================

@api_router.post("/startups/{startup_id}/milestones")
async def create_milestone(startup_id: str, body: MilestoneCreate, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    milestone = {
        "id": str(uuid.uuid4()),
        "startup_id": startup_id,
        "title": body.title,
        "description": body.description or "",
        "target_date": body.target_date if body.target_date else None,
        "status": "pending",
        "created_at": datetime.now(timezone.utc).isoformat(),
        "updated_at": datetime.now(timezone.utc).isoformat(),
    }
    supabase_client.table("milestones").insert(milestone).execute()
    return milestone

@api_router.get("/startups/{startup_id}/milestones")
async def get_milestones(startup_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    milestones_response = supabase_client.table("milestones").select("*").eq("startup_id", startup_id).execute()
    milestones = milestones_response.data
    
    for m in milestones:
        tasks_response = supabase_client.table("tasks").select("*").eq("milestone_id", m["id"]).execute()
        tasks = tasks_response.data
        
        total = len(tasks)
        done = len([t for t in tasks if t.get("status") == "done"])
        m["progress"] = int((done / total) * 100) if total > 0 else 0
        m["task_count"] = total
        m["tasks_done"] = done
        
    return milestones

@api_router.put("/milestones/{milestone_id}")
async def update_milestone(milestone_id: str, body: MilestoneUpdate, user=Depends(get_current_user)):
    milestone_response = supabase_client.table("milestones").select("*").eq("id", milestone_id).execute()
    milestone = milestone_response.data[0] if milestone_response.data else None
    
    if not milestone:
        raise HTTPException(status_code=404, detail="Milestone not found")
        
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", milestone["startup_id"]).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    updates = {"updated_at": datetime.now(timezone.utc).isoformat()}
    for field in ["title", "description", "target_date", "status"]:
        val = getattr(body, field, None)
        if field == "target_date" and val == "":
             updates[field] = None
        elif val is not None:
            updates[field] = val
            
    response = supabase_client.table("milestones").update(updates).eq("id", milestone_id).execute()
    return response.data[0] if response.data else None

@api_router.delete("/milestones/{milestone_id}")
async def delete_milestone(milestone_id: str, user=Depends(get_current_user)):
    milestone_response = supabase_client.table("milestones").select("*").eq("id", milestone_id).execute()
    milestone = milestone_response.data[0] if milestone_response.data else None
    
    if not milestone:
        raise HTTPException(status_code=404, detail="Milestone not found")
        
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", milestone["startup_id"]).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    supabase_client.table("milestones").delete().eq("id", milestone_id).execute()
    supabase_client.table("tasks").update({"milestone_id": None}).eq("milestone_id", milestone_id).execute()
    
    return {"success": True}

# ==================== FEEDBACK ROUTES ====================

@api_router.post("/startups/{startup_id}/feedback")
async def create_feedback(startup_id: str, body: FeedbackCreate, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    feedback = {
        "id": str(uuid.uuid4()),
        "startup_id": startup_id,
        "title": body.title,
        "content": body.content or "",
        "category": body.category or "product",
        "rating": body.rating or 3,
        "submitted_by": str(user.id),
        "source": body.source or "internal",
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    
    supabase_client.table("feedback").insert(feedback).execute()
    return feedback

@api_router.get("/startups/{startup_id}/feedback")
async def get_feedback(startup_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    response = supabase_client.table("feedback").select("*").eq("startup_id", startup_id).execute()
    return response.data

# ==================== ANALYTICS ROUTES ====================

@api_router.get("/startups/{startup_id}/analytics")
async def get_analytics(startup_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    tasks = supabase_client.table("tasks").select("*").eq("startup_id", startup_id).execute().data
    milestones = supabase_client.table("milestones").select("*").eq("startup_id", startup_id).execute().data
    feedbacks = supabase_client.table("feedback").select("*").eq("startup_id", startup_id).execute().data
    members = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).execute().data

    task_stats = {"todo": 0, "in_progress": 0, "review": 0, "done": 0}
    priority_stats = {"low": 0, "medium": 0, "high": 0, "urgent": 0}
    for t in tasks:
        status = t.get("status", "todo")
        if status in task_stats:
            task_stats[status] += 1
        prio = t.get("priority", "medium")
        if prio in priority_stats:
            priority_stats[prio] += 1

    milestone_stats = {"pending": 0, "in_progress": 0, "completed": 0}
    for m in milestones:
        ms = m.get("status", "pending")
        if ms in milestone_stats:
            milestone_stats[ms] += 1

    feedback_by_category = {}
    avg_rating = 0
    if feedbacks:
        for f in feedbacks:
            cat = f.get("category", "other")
            feedback_by_category[cat] = feedback_by_category.get(cat, 0) + 1
        avg_rating = round(sum(f.get("rating", 0) for f in feedbacks) / len(feedbacks), 1)

    total_tasks = len(tasks)
    completed_tasks = task_stats["done"]
    completion_rate = round((completed_tasks / total_tasks) * 100) if total_tasks > 0 else 0

    return {
        "total_tasks": total_tasks,
        "completed_tasks": completed_tasks,
        "completion_rate": completion_rate,
        "task_stats": task_stats,
        "priority_stats": priority_stats,
        "total_milestones": len(milestones),
        "milestone_stats": milestone_stats,
        "total_feedback": len(feedbacks),
        "feedback_by_category": feedback_by_category,
        "avg_rating": avg_rating,
        "team_size": len(members),
    }

# ==================== AI ROUTES (GEMINI) ====================

@api_router.post("/ai/insights")
async def get_ai_insights(body: AIInsightRequest, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", body.startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    startup_response = supabase_client.table("startups").select("*").eq("id", body.startup_id).execute()
    startup = startup_response.data[0] if startup_response.data else {}

    tasks = supabase_client.table("tasks").select("*").eq("startup_id", body.startup_id).execute().data
    milestones = supabase_client.table("milestones").select("*").eq("startup_id", body.startup_id).execute().data
    feedbacks = supabase_client.table("feedback").select("*").eq("startup_id", body.startup_id).execute().data

    task_summary = f"Total tasks: {len(tasks)}, Done: {len([t for t in tasks if t.get('status')=='done'])}, In Progress: {len([t for t in tasks if t.get('status')=='in_progress'])}"
    milestone_summary = f"Total milestones: {len(milestones)}, Completed: {len([m for m in milestones if m.get('status')=='completed'])}"
    feedback_summary = f"Total feedback: {len(feedbacks)}"
    if feedbacks:
        avg = round(sum(f.get("rating", 0) for f in feedbacks) / len(feedbacks), 1)
        feedback_summary += f", Average rating: {avg}/5"

    prompt_map = {
        "general": f"Analyze this startup's progress and provide 3-5 actionable insights:\nStartup: {startup.get('name', 'Unknown')} ({startup.get('industry', 'Unknown')} - {startup.get('stage', 'idea')} stage)\n{task_summary}\n{milestone_summary}\n{feedback_summary}\nProvide specific, actionable recommendations.",
        "tasks": f"Suggest 5 strategic tasks for this startup:\nStartup: {startup.get('name', 'Unknown')} in {startup.get('industry', 'Unknown')} at {startup.get('stage', 'idea')} stage.\nCurrent: {task_summary}\nProvide with title, description, and priority.",
        "milestones": f"Suggest 3 key milestones:\nStartup: {startup.get('name', 'Unknown')} in {startup.get('industry', 'Unknown')} at {startup.get('stage', 'idea')} stage.\nCurrent: {milestone_summary}",
        "growth": f"Growth strategy analysis:\nStartup: {startup.get('name', 'Unknown')} - {startup.get('industry', 'Unknown')}\n{task_summary}\n{milestone_summary}\n{feedback_summary}"
    }
    prompt = prompt_map.get(body.prompt_type, prompt_map["general"])
    
    # Add custom prompt if provided
    if body.custom_prompt:
        prompt += f"\n\nAdditional context from user: {body.custom_prompt}"

    try:
        import google.generativeai as genai
        genai.configure(api_key=GEMINI_API_KEY)
        model = genai.GenerativeModel(
            model_name="gemini-2.5-flash",
            system_instruction="You are a startup advisor. Provide concise, actionable insights."
        )
        response = model.generate_content(prompt)
        insights_text = response.text
        
        # Save to history
        try:
            supabase_client.table("ai_history").insert({
                "startup_id": body.startup_id,
                "type": "insight",
                "subtype": body.prompt_type,
                "content": insights_text,
                "metadata": {"prompt_type": body.prompt_type},
                "created_by": str(user.id)
            }).execute()
        except Exception as hist_err:
            logger.warning(f"Failed to save insight to history: {hist_err}")
        
        return {"insights": insights_text, "prompt_type": body.prompt_type}
    except Exception as e:
        logger.error(f"AI insights error: {e}")
        raise HTTPException(status_code=500, detail=f"AI service error: {str(e)}")

@api_router.post("/ai/pitch")
async def generate_pitch(body: PitchRequest, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", body.startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    startup_response = supabase_client.table("startups").select("*").eq("id", body.startup_id).execute()
    startup = startup_response.data[0] if startup_response.data else {}

    tasks = supabase_client.table("tasks").select("*").eq("startup_id", body.startup_id).execute().data
    milestones = supabase_client.table("milestones").select("*").eq("startup_id", body.startup_id).execute().data
    feedbacks = supabase_client.table("feedback").select("*").eq("startup_id", body.startup_id).execute().data
    members = supabase_client.table("startup_members").select("*").eq("startup_id", body.startup_id).execute().data

    completed_tasks = len([t for t in tasks if t.get("status") == "done"])
    completed_milestones = len([m for m in milestones if m.get("status") == "completed"])
    avg_rating = round(sum(f.get("rating", 0) for f in feedbacks) / len(feedbacks), 1) if feedbacks else 0

    prompt = f"""Generate a professional investor pitch as a structured set of PowerPoint slides in JSON format. 

Company: {startup.get('name', 'Unknown')}
Industry: {startup.get('industry', 'Unknown')}
Stage: {startup.get('stage', 'idea')}
Description: {startup.get('description', 'No description')}

Traction: Team {len(members)}, Tasks {completed_tasks}/{len(tasks)}, Milestones {completed_milestones}/{len(milestones)}, Rating {avg_rating}/5

Return ONLY a valid JSON object with this structure:
{{
  "title": "Pitch Deck",
  "slides": [
    {{"title": "Title Slide", "content": ["Company Name", "Tagline", "Contact"]}},
    {{"title": "Problem", "content": ["Key problem statement", "Market pain point", "Why it matters"]}},
    {{"title": "Solution", "content": ["Our solution", "Key features", "How it solves the problem"]}},
    {{"title": "Market Opportunity", "content": ["TAM/SAM/SOM", "Market size", "Growth potential"]}},
    {{"title": "Traction", "content": ["Milestones: {completed_milestones}/{len(milestones)}", "Tasks completed: {completed_tasks}/{len(tasks)}", "Team size: {len(members)}"]}},
    {{"title": "Business Model", "content": ["Revenue streams", "Pricing strategy", "Unit economics"]}},
    {{"title": "Team", "content": ["Key team members", "Expertise", "Track record"]}},
    {{"title": "Roadmap", "content": ["3-6 month goals", "12 month vision", "Key milestones"]}},
    {{"title": "Financial Projections", "content": ["Year 1 revenue", "Year 2 revenue", "Path to profitability"]}},
    {{"title": "The Ask", "content": ["Funding amount", "Use of funds", "Expected outcomes"]}}
  ]
}}"""
    
    # Add custom prompt if provided
    if body.custom_prompt:
        prompt += f"\n\nAdditional context from user: {body.custom_prompt}\nPlease incorporate this context into the pitch generation."

    try:
        import google.generativeai as genai
        genai.configure(api_key=GEMINI_API_KEY)
        model = genai.GenerativeModel(
            model_name="gemini-2.5-flash",
            system_instruction="You are a professional pitch deck creator. Generate ONLY valid JSON output with no markdown, no code blocks, no extra text. Output must start with { and end with }"
        )
        response = model.generate_content(prompt)
        
        # Clean up response - remove markdown code blocks if present
        response_text = response.text.strip()
        if response_text.startswith("```json"):
            response_text = response_text[7:]
        if response_text.startswith("```"):
            response_text = response_text[3:]
        if response_text.endswith("```"):
            response_text = response_text[:-3]
        response_text = response_text.strip()
        
        # Attempt to parse the JSON response with error recovery
        try:
            pitch_json = json.loads(response_text)
        except json.JSONDecodeError as initial_error:
            # Try to fix common JSON issues
            logger.warning(f"Initial JSON parse failed, attempting repair...")
            
            # Replace problematic characters in content strings
            # Fix newlines within string values (replace \n with space in content arrays)
            fixed_text = response_text
            
            # More aggressive: find all content arrays and fix them
            import re
            
            # Find all "content": [...] patterns and fix newlines inside
            def fix_content_array(match):
                content_section = match.group(0)
                # Replace literal newlines with spaces within the array
                content_section = content_section.replace('\n', ' ').replace('\r', ' ')
                return content_section
            
            fixed_text = re.sub(r'"content":\s*\[[^\]]*\]', fix_content_array, fixed_text, flags=re.DOTALL)
            
            try:
                pitch_json = json.loads(fixed_text)
            except json.JSONDecodeError as second_error:
                # Last resort: create default slides
                logger.error(f"JSON repair failed: {second_error}")
                logger.error(f"Response text: {response_text[:1000]}")
                
                # Return a default pitch structure
                pitch_json = {
                    "title": "Pitch Deck",
                    "slides": [
                        {"title": "Title Slide", "content": [startup.get('name', 'Startup'), startup.get('industry', 'Tech'), "Your Company"]},
                        {"title": "Problem", "content": [startup.get('description', 'Solving a market problem'), "Key market need"]},
                        {"title": "Solution", "content": ["Our innovative solution", "Key features and benefits"]},
                        {"title": "Market Opportunity", "content": [f"Large {startup.get('industry', 'market')} opportunity", "Growth potential"]},
                        {"title": "The Ask", "content": ["Seeking investment", "To accelerate growth"]}
                    ]
                }
        
        # Save to history
        pitch_content = json.dumps(pitch_json)
        try:
            supabase_client.table("ai_history").insert({
                "startup_id": body.startup_id,
                "type": "pitch",
                "subtype": "pitch_deck",
                "content": pitch_content,
                "metadata": {"slide_count": len(pitch_json.get("slides", []))},
                "created_by": str(user.id)
            }).execute()
        except Exception as hist_err:
            logger.warning(f"Failed to save pitch to history: {hist_err}")
        
        return {
            "pitch": pitch_content,
            "startup_name": startup.get("name", ""),
            "slides": pitch_json.get("slides", []),
            "format": "ppt"
        }
    except Exception as e:
        logger.error(f"Pitch error: {e}")
        raise HTTPException(status_code=500, detail=f"AI service error: {str(e)}")

@api_router.post("/ai/pitch/download")
async def download_pitch_pptx(body: PitchRequest, user=Depends(get_current_user)):
    """Generate and download pitch as PowerPoint file"""
    startup_response = supabase_client.table("startups").select("*").eq("id", body.startup_id).execute()
    startup = startup_response.data[0] if startup_response.data else {}

    tasks = supabase_client.table("tasks").select("*").eq("startup_id", body.startup_id).execute().data
    milestones = supabase_client.table("milestones").select("*").eq("startup_id", body.startup_id).execute().data
    feedbacks = supabase_client.table("feedback").select("*").eq("startup_id", body.startup_id).execute().data
    members = supabase_client.table("startup_members").select("*").eq("startup_id", body.startup_id).execute().data

    completed_tasks = len([t for t in tasks if t.get("status") == "done"])
    completed_milestones = len([m for m in milestones if m.get("status") == "completed"])

    prompt = f"""Generate a professional investor pitch as a structured set of PowerPoint slides in JSON format.
Company: {startup.get('name', 'Unknown')}
Industry: {startup.get('industry', 'Unknown')}
Stage: {startup.get('stage', 'idea')}
Description: {startup.get('description', 'No description')}
Traction: Team {len(members)}, Tasks {completed_tasks}/{len(tasks)}, Milestones {completed_milestones}/{len(milestones)}

Return ONLY a valid JSON object with this structure:
{{
  "slides": [
    {{"title": "Title Slide", "content": ["Company Name", "Tagline", "Contact"]}},
    {{"title": "Problem", "content": ["Key problem statement", "Market pain point"]}},
    {{"title": "Solution", "content": ["Our solution", "Key features"]}},
    {{"title": "Market Opportunity", "content": ["TAM/SAM/SOM", "Market size"]}},
    {{"title": "Traction", "content": ["Milestones: {completed_milestones}/{len(milestones)}", "Tasks completed: {completed_tasks}/{len(tasks)}", "Team size: {len(members)}"]}},
    {{"title": "Business Model", "content": ["Revenue streams", "Pricing strategy"]}},
    {{"title": "Team", "content": ["Key team members", "Expertise"]}},
    {{"title": "Roadmap", "content": ["3-6 month goals", "12 month vision"]}},
    {{"title": "The Ask", "content": ["Funding amount", "Use of funds"]}}
  ]
}}"""

    try:
        import google.generativeai as genai
        genai.configure(api_key=GEMINI_API_KEY)
        model = genai.GenerativeModel(
            model_name="gemini-2.5-flash",
            system_instruction="You are a professional pitch deck creator. Generate ONLY valid JSON output with no markdown, no code blocks, no extra text. Output must start with { and end with }"
        )
        response = model.generate_content(prompt)
        
        # Clean up response - remove markdown code blocks if present
        response_text = response.text.strip()
        if response_text.startswith("```json"):
            response_text = response_text[7:]
        if response_text.startswith("```"):
            response_text = response_text[3:]
        if response_text.endswith("```"):
            response_text = response_text[:-3]
        response_text = response_text.strip()
        
        pitch_data = json.loads(response_text)
        
        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        for slide_info in pitch_data.get("slides", []):
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = slide_info.get("title", "Slide")
            title_frame.paragraphs[0].font.size = Pt(54)
            title_frame.paragraphs[0].font.bold = True
            
            # Add content
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for i, bullet in enumerate(slide_info.get("content", [])):
                if i == 0:
                    content_frame.text = f"• {bullet}"
                else:
                    p = content_frame.add_paragraph()
                    p.text = f"• {bullet}"
                    p.level = 0
                p = content_frame.paragraphs[i]
                p.font.size = Pt(28)
        
        # Save to bytes
        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        startup_name = startup.get("name", "pitch").replace(" ", "_").lower()
        return StreamingResponse(
            iter([pptx_bytes.getvalue()]),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={startup_name}_pitch.pptx"}
        )
    except json.JSONDecodeError as e:
        logger.error(f"Download: Failed to parse JSON: {response.text[:500]}")
        raise HTTPException(status_code=500, detail=f"Error generating presentation: Invalid JSON")
    except Exception as e:
        logger.error(f"PPTX generation error: {e}")
        raise HTTPException(status_code=500, detail=f"Error generating presentation: {str(e)}")

# ==================== AI HISTORY ROUTES ====================

@api_router.get("/ai/history/{startup_id}")
async def get_ai_history(startup_id: str, ai_type: str = None, user=Depends(get_current_user)):
    """Get AI generation history for a startup"""
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
    
    query = supabase_client.table("ai_history").select("*").eq("startup_id", startup_id)
    if ai_type:
        query = query.eq("type", ai_type)
    
    response = query.order("created_at", desc=True).limit(50).execute()
    history = []
    
    for item in response.data:
        history.append({
            "id": item["id"],
            "type": item["type"],
            "subtype": item["subtype"],
            "created_at": item["created_at"],
            "preview": item["content"][:100] + "..." if len(item["content"]) > 100 else item["content"],
            "metadata": item.get("metadata", {})
        })
    
    return {"history": history}

@api_router.get("/ai/history/{startup_id}/{history_id}")
async def get_history_item(startup_id: str, history_id: str, user=Depends(get_current_user)):
    """Get a specific history item"""
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
    
    response = supabase_client.table("ai_history").select("*").eq("id", history_id).eq("startup_id", startup_id).execute()
    if not response.data:
        raise HTTPException(status_code=404, detail="History item not found")
    
    item = response.data[0]
    return {
        "id": item["id"],
        "type": item["type"],
        "subtype": item["subtype"],
        "content": item["content"],
        "metadata": item.get("metadata", {}),
        "created_at": item["created_at"]
    }

@api_router.post("/ai/history")
async def save_ai_history(body: dict, user=Depends(get_current_user)):
    """Save AI generation to history"""
    startup_id = body.get("startup_id")
    ai_type = body.get("type")  # "insight" or "pitch"
    subtype = body.get("subtype")  # e.g., "general", "pitch_deck", etc.
    content = body.get("content")
    metadata = body.get("metadata", {})
    
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
    
    try:
        response = supabase_client.table("ai_history").insert({
            "startup_id": startup_id,
            "type": ai_type,
            "subtype": subtype,
            "content": content,
            "metadata": metadata,
            "created_by": str(user.id)
        }).execute()
        
        if response.data:
            return {"id": response.data[0]["id"], "message": "History saved"}
        else:
            raise HTTPException(status_code=500, detail="Failed to save history")
    except Exception as e:
        logger.error(f"Error saving history: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.delete("/ai/history/{history_id}")
async def delete_history(history_id: str, user=Depends(get_current_user)):
    """Delete a history item"""
    # Check ownership
    response = supabase_client.table("ai_history").select("*").eq("id", history_id).execute()
    if not response.data:
        raise HTTPException(status_code=404, detail="History item not found")
    
    item = response.data[0]
    
    # Verify user is a member of the startup
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", item["startup_id"]).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
    
    supabase_client.table("ai_history").delete().eq("id", history_id).execute()
    return {"message": "History deleted"}

# ==================== TEAM ROUTES ====================

@api_router.get("/startups/{startup_id}/members")
async def get_members(startup_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    members = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).execute().data
    result = []
    
    member_user_ids = [m["user_id"] for m in members]
    profiles_response = supabase_client.table("profiles").select("*").in_("id", member_user_ids).execute()
    profiles_map = {p["id"]: p for p in profiles_response.data}
    
    for m in members:
        profile = profiles_map.get(m["user_id"])
        result.append({
            "id": m["id"],
            "user_id": m["user_id"],
            "role": m["role"],
            "joined_at": m["joined_at"],
            "email": profile.get("email", "") if profile else "",
            "full_name": profile.get("full_name", "") if profile else "",
            "avatar_url": profile.get("avatar_url", "") if profile else "",
        })
    return result

@api_router.delete("/startups/{startup_id}/members/{user_id}")
async def remove_member(startup_id: str, user_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    requester = mem_response.data[0] if mem_response.data else None
    
    if not requester or requester["role"] != "founder":
        raise HTTPException(status_code=403, detail="Only founders can remove members")
        
    if user_id == str(user.id):
        raise HTTPException(status_code=400, detail="Cannot remove yourself")
        
    supabase_client.table("startup_members").delete().eq("startup_id", startup_id).eq("user_id", user_id).execute()
    return {"success": True}

@api_router.put("/startups/{startup_id}/members/{user_id}/role")
async def update_member_role(startup_id: str, user_id: str, body: MemberRoleUpdate, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    requester = mem_response.data[0] if mem_response.data else None
    
    if not requester or requester["role"] != "founder":
        raise HTTPException(status_code=403, detail="Only founders can change roles")
    
    target_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", user_id).execute()
    target_member = target_response.data[0] if target_response.data else None

    if not target_member:
        raise HTTPException(status_code=404, detail="Member not found")
    
    if target_member["role"] == "founder":
        raise HTTPException(status_code=400, detail="Cannot change founder role")
    
    if body.role not in ["manager", "member", "investor"]:
        raise HTTPException(status_code=400, detail="Invalid role.")
    
    supabase_client.table("startup_members").update({"role": body.role}).eq("startup_id", startup_id).eq("user_id", user_id).execute()
    return {"success": True, "role": body.role}

@api_router.get("/startups/{startup_id}/invite-code")
async def get_invite_code(startup_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    member = mem_response.data[0] if mem_response.data else None
    
    if not member or member["role"] != "founder":
        raise HTTPException(status_code=403, detail="Only founders can view invite code")
        
    startup_response = supabase_client.table("startups").select("*").eq("id", startup_id).execute()
    startup = startup_response.data[0] if startup_response.data else {}
    return {"invite_code": startup.get("invite_code", "")}


@api_router.get("/startups/{startup_id}/investor-view")
async def get_investor_view(startup_id: str, user=Depends(get_current_user)):
    """Get startup metrics and performance data for investors."""
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
    
    # Get startup details
    startup_res = supabase_client.table("startups").select("*").eq("id", startup_id).execute()
    startup = startup_res.data[0] if startup_res.data else {}
    
    # Get metrics
    tasks_res = supabase_client.table("tasks").select("status").eq("startup_id", startup_id).execute()
    milestones_res = supabase_client.table("milestones").select("status").eq("startup_id", startup_id).execute()
    income_res = supabase_client.table("income").select("amount").eq("startup_id", startup_id).execute()
    expenses_res = supabase_client.table("expenses").select("amount").eq("startup_id", startup_id).execute()
    investments_res = supabase_client.table("investments").select("amount").eq("startup_id", startup_id).execute()
    members_res = supabase_client.table("startup_members").select("id").eq("startup_id", startup_id).execute()
    
    tasks_completed = len([t for t in tasks_res.data if t.get("status") == "done"]) if tasks_res.data else 0
    tasks_total = len(tasks_res.data) if tasks_res.data else 0
    
    milestones_completed = len([m for m in milestones_res.data if m.get("status") == "completed"]) if milestones_res.data else 0
    milestones_total = len(milestones_res.data) if milestones_res.data else 0
    
    total_income = sum(inc.get("amount", 0) for inc in income_res.data) if income_res.data else 0
    total_expenses = sum(exp.get("amount", 0) for exp in expenses_res.data) if expenses_res.data else 0
    total_raised = sum(inv.get("amount", 0) for inv in investments_res.data) if investments_res.data else 0
    team_size = len(members_res.data) if members_res.data else 0
    
    # Calculate current balance
    current_balance = total_income + total_raised - total_expenses
    
    # Calculate monthly burn rate (average expenses per month)
    # This is a simplified calculation - in production you'd want to group by month
    monthly_burn = total_expenses / 12 if total_expenses > 0 else 0
    
    # Calculate runway (months of operations possible)
    runway = current_balance / monthly_burn if monthly_burn > 0 else 0
    
    return {
        "startup": startup,
        "metrics": {
            "tasks_completed": tasks_completed,
            "tasks_total": tasks_total,
            "milestones_completed": milestones_completed,
            "milestones_total": milestones_total,
            "total_income": total_income,
            "total_expenses": total_expenses,
            "total_raised": total_raised,
            "current_balance": current_balance,
            "team_size": team_size,
            "monthly_burn": monthly_burn,
            "runway": runway,
        }
    }


@api_router.get("/startups/{startup_id}/investors")
async def get_investors(startup_id: str, user=Depends(get_current_user)):
    """Get list of investors and pending invites for a startup."""
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    member = mem_response.data[0] if mem_response.data else None
    
    if not member or member["role"] != "founder":
        raise HTTPException(status_code=403, detail="Only founders can view investors")
    
    # Get investor members
    investors_res = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("role", "investor").execute()
    investors_list = []
    
    for inv in investors_res.data or []:
        profile_res = supabase_client.table("profiles").select("*").eq("id", inv["user_id"]).execute()
        profile = profile_res.data[0] if profile_res.data else {}
        investors_list.append({
            "id": inv["id"],
            "user_id": inv["user_id"],
            "full_name": profile.get("full_name", ""),
            "email": profile.get("email", ""),
            "joined_at": inv["joined_at"],
        })
    
    return {
        "investors": investors_list,
        "pending_invites": []
    }


@api_router.post("/startups/{startup_id}/investors/invite")
async def invite_investor(startup_id: str, body: dict, user=Depends(get_current_user)):
    """Send invite to an investor."""
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    member = mem_response.data[0] if mem_response.data else None
    
    if not member or member["role"] != "founder":
        raise HTTPException(status_code=403, detail="Only founders can invite investors")
    
    # Return the startup's invite code
    startup_res = supabase_client.table("startups").select("invite_code").eq("id", startup_id).execute()
    startup = startup_res.data[0] if startup_res.data else {}
    
    return {"invite_code": startup.get("invite_code", "")}


@api_router.delete("/startups/{startup_id}/investors/{user_id}")
async def remove_investor(startup_id: str, user_id: str, user=Depends(get_current_user)):
    """Remove an investor from a startup."""
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    requester = mem_response.data[0] if mem_response.data else None
    
    if not requester or requester["role"] != "founder":
        raise HTTPException(status_code=403, detail="Only founders can remove investors")
    
    supabase_client.table("startup_members").delete().eq("startup_id", startup_id).eq("user_id", user_id).eq("role", "investor").execute()
    return {"success": True}

@api_router.post("/startups/{startup_id}/regenerate-invite")
async def regenerate_invite(startup_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    member = mem_response.data[0] if mem_response.data else None
    
    if not member or member["role"] != "founder":
        raise HTTPException(status_code=403, detail="Only founders can regenerate invite")
        
    new_code = str(uuid.uuid4())[:8].upper()
    supabase_client.table("startups").update({"invite_code": new_code}).eq("id", startup_id).execute()
    return {"invite_code": new_code}

# ==================== SUBSCRIPTION ROUTES ====================

@api_router.get("/startups/{startup_id}/subscription")
async def get_subscription(startup_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    sub_response = supabase_client.table("subscriptions").select("*").eq("startup_id", startup_id).execute()
    sub = sub_response.data[0] if sub_response.data else None
    
    if not sub:
        sub = {
            "id": str(uuid.uuid4()),
            "startup_id": startup_id,
            "plan": "free",
            "status": "active",
            "created_at": datetime.now(timezone.utc).isoformat(),
            "updated_at": datetime.now(timezone.utc).isoformat(),
        }
        supabase_client.table("subscriptions").insert(sub).execute()
        return sub
    return sub

@api_router.post("/startups/{startup_id}/subscription")
async def update_subscription(startup_id: str, body: SubscriptionUpdate, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    member = mem_response.data[0] if mem_response.data else None
    
    if not member or member["role"] != "founder":
        raise HTTPException(status_code=403, detail="Only founders can manage subscription")
        
    sub_response = supabase_client.table("subscriptions").select("*").eq("startup_id", startup_id).execute()
    existing = sub_response.data[0] if sub_response.data else None
    
    if existing:
        supabase_client.table("subscriptions").update(
            {"plan": body.plan, "status": "active", "updated_at": datetime.now(timezone.utc).isoformat()}
        ).eq("startup_id", startup_id).execute()
    else:
        sub = {
            "id": str(uuid.uuid4()),
            "startup_id": startup_id,
            "plan": body.plan,
            "status": "active",
            "created_at": datetime.now(timezone.utc).isoformat(),
            "updated_at": datetime.now(timezone.utc).isoformat(),
        }
        supabase_client.table("subscriptions").insert(sub).execute()
        
    supabase_client.table("startups").update({"subscription_plan": body.plan}).eq("id", startup_id).execute()
    
    response = supabase_client.table("subscriptions").select("*").eq("startup_id", startup_id).execute()
    return response.data[0]

# ==================== FINANCE ROUTES ====================

@api_router.post("/startups/{startup_id}/finance/income")
async def create_income(startup_id: str, body: IncomeCreate, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    member = mem_response.data[0] if mem_response.data else None
    
    if not member or member["role"] not in ["founder", "manager"]:
        raise HTTPException(status_code=403, detail="Founders/managers only")
        
    income = {
        "id": str(uuid.uuid4()),
        "startup_id": startup_id,
        "title": body.title,
        "amount": body.amount,
        "category": body.category or "revenue",
        "date": body.date or datetime.now(timezone.utc).strftime("%Y-%m-%d"),
        "notes": body.notes or "",
        "created_by": str(user.id),
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    
    supabase_client.table("income").insert(income).execute()
    return income

@api_router.get("/startups/{startup_id}/finance/income")
async def get_income(startup_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    response = supabase_client.table("income").select("*").eq("startup_id", startup_id).order("date", desc=True).limit(500).execute()
    return response.data

@api_router.delete("/startups/{startup_id}/finance/income/{income_id}")
async def delete_income(startup_id: str, income_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    member = mem_response.data[0] if mem_response.data else None
    
    if not member or member["role"] not in ["founder", "manager"]:
        raise HTTPException(status_code=403, detail="Founders/managers only")
        
    supabase_client.table("income").delete().eq("id", income_id).execute()
    return {"success": True}

@api_router.post("/startups/{startup_id}/finance/expenses")
async def create_expense(startup_id: str, body: ExpenseCreate, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    member = mem_response.data[0] if mem_response.data else None
    
    if not member or member["role"] not in ["founder", "manager"]:
        raise HTTPException(status_code=403, detail="Founders/managers only")
        
    expense = {
        "id": str(uuid.uuid4()),
        "startup_id": startup_id,
        "title": body.title,
        "amount": body.amount,
        "category": body.category or "operations",
        "date": body.date or datetime.now(timezone.utc).strftime("%Y-%m-%d"),
        "notes": body.notes or "",
        "created_by": str(user.id),
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    
    supabase_client.table("expenses").insert(expense).execute()
    return expense

@api_router.get("/startups/{startup_id}/finance/expenses")
async def get_expenses(startup_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    response = supabase_client.table("expenses").select("*").eq("startup_id", startup_id).order("date", desc=True).limit(500).execute()
    return response.data

@api_router.delete("/startups/{startup_id}/finance/expenses/{expense_id}")
async def delete_expense(startup_id: str, expense_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    member = mem_response.data[0] if mem_response.data else None
    
    if not member or member["role"] not in ["founder", "manager"]:
        raise HTTPException(status_code=403, detail="Founders/managers only")
        
    supabase_client.table("expenses").delete().eq("id", expense_id).execute()
    return {"success": True}

@api_router.post("/startups/{startup_id}/finance/investments")
async def create_investment(startup_id: str, body: InvestmentCreate, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    member = mem_response.data[0] if mem_response.data else None
    
    if not member:
        raise HTTPException(status_code=403, detail="Not a member of this startup")
    
    # Allow both founders and investors to create investment records
    if member["role"] not in ["founder", "investor"]:
        raise HTTPException(status_code=403, detail="Only founders and investors can add investments")
        
    investment = {
        "id": str(uuid.uuid4()),
        "startup_id": startup_id,
        "investor_name": body.investor_name,
        "amount": body.amount,
        "equity_percentage": body.equity_percentage or 0,
        "investment_type": body.investment_type or "seed",
        "date": body.date or datetime.now(timezone.utc).strftime("%Y-%m-%d"),
        "notes": body.notes or "",
        "created_by": str(user.id),
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    supabase_client.table("investments").insert(investment).execute()
    return investment

@api_router.get("/startups/{startup_id}/finance/investments")
async def get_investments(startup_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    response = supabase_client.table("investments").select("*").eq("startup_id", startup_id).order("date", desc=True).limit(100).execute()
    return response.data

@api_router.delete("/startups/{startup_id}/finance/investments/{investment_id}")
async def delete_investment(startup_id: str, investment_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    member = mem_response.data[0] if mem_response.data else None
    
    if not member:
        raise HTTPException(status_code=403, detail="Not a member of this startup")
    
    # Allow founders and investors to delete investments
    if member["role"] not in ["founder", "investor"]:
        raise HTTPException(status_code=403, detail="Only founders and investors can delete investments")
        
    supabase_client.table("investments").delete().eq("id", investment_id).execute()
    return {"success": True}

@api_router.get("/startups/{startup_id}/finance/summary")
async def get_finance_summary(startup_id: str, user=Depends(get_current_user)):
    mem_response = supabase_client.table("startup_members").select("*").eq("startup_id", startup_id).eq("user_id", str(user.id)).execute()
    if not mem_response.data:
        raise HTTPException(status_code=403, detail="Not a member")
        
    income = supabase_client.table("income").select("*").eq("startup_id", startup_id).execute().data
    expenses = supabase_client.table("expenses").select("*").eq("startup_id", startup_id).execute().data
    investments = supabase_client.table("investments").select("*").eq("startup_id", startup_id).execute().data
    
    total_income = sum(i.get("amount", 0) for i in income)
    total_expenses = sum(e.get("amount", 0) for e in expenses)
    total_investments = sum(inv.get("amount", 0) for inv in investments)
    net_balance = total_income + total_investments - total_expenses
    
    # Calculate runway
    monthly_burn = total_expenses / 12 if total_expenses > 0 else 0
    runway_months = net_balance / monthly_burn if monthly_burn > 0 else 0
    
    # Calculate total equity given
    total_equity_given = sum(inv.get("equity_percentage", 0) for inv in investments)
    
    # Calculate income and expenses by category
    income_by_category = {}
    for item in income:
        cat = item.get("category", "other")
        income_by_category[cat] = income_by_category.get(cat, 0) + item.get("amount", 0)
    
    expenses_by_category = {}
    for item in expenses:
        cat = item.get("category", "other")
        expenses_by_category[cat] = expenses_by_category.get(cat, 0) + item.get("amount", 0)
    
    return {
        "total_income": total_income,
        "total_expenses": total_expenses,
        "total_investments": total_investments,
        "net_balance": net_balance,
        "runway_months": round(runway_months),
        "total_equity_given": total_equity_given,
        "income_by_category": income_by_category,
        "expenses_by_category": expenses_by_category,
    }

# ==================== INVESTOR ROUTES ====================

class SwipeAction(BaseModel):
    action: str  # 'interested' or 'passed'

@app.get("/api/investor/browse")
async def browse_startups(user: dict = Depends(get_current_user)):
    """
    Get a list of startups for the investor to swipe on.
    Excludes startups the user has already swiped on.
    """
    try:
        # 1. Get all swipes by this user
        swipes_res = supabase_client.table("investor_swipes").select("startup_id").eq("investor_id", str(user.id)).execute()
        swiped_ids = [s['startup_id'] for s in swipes_res.data]
        
        # 2. Get all startups
        startups_res = supabase_client.table("startups").select("*").execute()
        all_startups = startups_res.data
        
        # 3. Filter out swiped startups
        # Also fetching enriched data like milestones/finances could be done here or in separate calls
        # For this MVP, we return the base startup data.
        
        available = []
        for s in all_startups:
            if s['id'] in swiped_ids:
                continue
            
            # Enrich with some finance data for the card
            # (In production, do this with a join view)
            try:
                # Get total raised
                inv_res = supabase_client.table("investments").select("amount").eq("startup_id", s['id']).execute()
                total_raised = sum(float(i['amount']) for i in inv_res.data)
                s['total_raised'] = total_raised
                
                # Get milestones progress
                mil_res = supabase_client.table("milestones").select("status").eq("startup_id", s['id']).execute()
                total_milestones = len(mil_res.data)
                completed_milestones = len([m for m in mil_res.data if m['status'] == 'completed'])
                s['milestones_total'] = total_milestones
                s['milestones_completed'] = completed_milestones

                # Get team size
                team_res = supabase_client.table("startup_members").select("id", count="exact").eq("startup_id", s['id']).execute()
                s['team_size'] = team_res.count if team_res.count else 1 # +1 for founder
            except Exception as e:
                logger.error(f"Error enriching startup {s['id']}: {e}")

            available.append(s)
            
        return available
    except Exception as e:
        logger.error(f"Browse error: {e}")
        raise HTTPException(status_code=500, detail="Failed to fetch startups")

@app.post("/api/investor/swipe/{startup_id}")
async def swipe_startup(startup_id: str, body: SwipeAction, user: dict = Depends(get_current_user)):
    try:
        # Check if already swiped
        existing = supabase_client.table("investor_swipes").select("*").eq("investor_id", str(user.id)).eq("startup_id", startup_id).execute()
        if existing.data:
            return {"message": "Already swiped"}

        # Record swipe
        data = {
            "investor_id": str(user.id),
            "startup_id": startup_id,
            "action": body.action,
            "created_at": datetime.now(timezone.utc).isoformat()
        }
        supabase_client.table("investor_swipes").insert(data).execute()
        return {"message": f"Recorded {body.action}"}
    except Exception as e:
        logger.error(f"Swipe error: {e}")
        raise HTTPException(status_code=500, detail="Failed to record swipe")

@app.get("/api/investor/matches")
async def get_matches(user: dict = Depends(get_current_user)):
    """
    Get all startups the investor swiped right on.
    """
    try:
        # Get 'interested' swipes
        swipes_res = supabase_client.table("investor_swipes").select("startup_id, created_at").eq("investor_id", str(user.id)).eq("action", "interested").execute()
        
        matches = []
        for swipe in swipes_res.data:
            startup_id = swipe['startup_id']
            # Fetch startup details
            s_res = supabase_client.table("startups").select("*").eq("id", startup_id).execute()
            if s_res.data:
                startup = s_res.data[0]
                startup['swiped_at'] = swipe['created_at']
                
                # Fetch Founder contact info (assuming creator of startup is contact)
                # In real app, `created_by` in startups table points to founder
                founder_id = startup.get('created_by')
                if founder_id:
                    f_res = supabase_client.table("profiles").select("full_name, email").eq("id", founder_id).execute()
                    if f_res.data:
                        startup['founder_name'] = f_res.data[0]['full_name']
                        startup['founder_email'] = f_res.data[0]['email']
                
                matches.append(startup)
                
        return matches
    except Exception as e:
        logger.error(f"Matches error: {e}")


@app.delete("/api/investor/matches/{startup_id}")
async def remove_match(startup_id: str, user: dict = Depends(get_current_user)):
    """
    Remove a match by deleting the swipe record.
    """
    try:
        # Delete the 'interested' swipe for this startup
        supabase_client.table("investor_swipes").delete().eq("investor_id", str(user.id)).eq("startup_id", startup_id).eq("action", "interested").execute()
        return {"status": "success", "message": "Match removed"}
    except Exception as e:
        logger.error(f"Remove match error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


app.include_router(api_router)
